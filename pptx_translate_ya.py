'''
Created on 2018/09/22
Simple script for translate *.pptx file using request & python-pptx packages
@author: miro
'''
import requests
import json
from pptx import Presentation

api_key = 'input_api_key_here' # https://translate.yandex.com/developers

def translate(srt_in):
	url = 'https://translate.yandex.net/api/v1.5/tr.json/translate'
	payload = {'key': api_key,
			   'text': srt_in,
			   'lang': 'en-ru',
			   }
	r = requests.get(url, params=payload)
	content = r.text  # выводим текст ответа
	m = json.loads(content)  # парсим json
	str_out = ''.join(m['text'])  # преобразуем список в строчку
	return str_out

if __name__ == '__main__':
    path_to_presentation = "Presentation.pptx"
    name = path_to_presentation[:-5]

    prs = Presentation(path_to_presentation)

    print("START...")
    for ns, slide in enumerate(prs.slides):
        print(slide)
        for nsh, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for np, paragraph in enumerate(shape.text_frame.paragraphs):
                for rs, run in enumerate(paragraph.runs):
                    str_in = run.text
                    str_out = translate(str_in)
                    prs.slides[ns].shapes[nsh].text_frame.paragraphs[np].runs[rs].text = str_out
                    
    prs.save('{0}{1}'.format(name, '_rus.pptx'))
    print("...DONE for {}".format(path_to_presentation))
